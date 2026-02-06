import streamlit as st
from groq import Groq
import os

# ---------- File Readers ----------
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract

# ---------- Page Config ----------
st.set_page_config(
    page_title="File Summarizer",
    layout="centered"
)

st.title("📂File Summarizer")
st.write("Upload a file and choose how you want the summary.")

# ---------- Groq Client ----------
client = Groq(
    api_key=os.getenv("GROQ_API_KEY"),
    timeout=90
)

MODEL = "llama-3.1-8b-instant"

# ---------- Text Extraction ----------
def extract_text(file):
    name = file.name.lower()

    if name.endswith(".pdf"):
        reader = PdfReader(file)
        return " ".join(page.extract_text() or "" for page in reader.pages)

    elif name.endswith(".docx"):
        doc = Document(file)
        return " ".join(p.text for p in doc.paragraphs)

    elif name.endswith(".txt"):
        return file.read().decode("utf-8")

    elif name.endswith(".csv"):
        return pd.read_csv(file).to_string()

    elif name.endswith((".xls", ".xlsx")):
        return pd.read_excel(file).to_string()

    elif name.endswith(".pptx"):
        prs = Presentation(file)
        return " ".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, "text")
        )

    elif name.endswith((".html", ".htm")):
        return BeautifulSoup(file.read(), "html.parser").get_text()

    elif name.endswith((".png", ".jpg", ".jpeg")):
        return pytesseract.image_to_string(Image.open(file))

    return None

# ---------- Chunking ----------
def chunk_text(text, max_chars=3000):
    return [text[i:i + max_chars] for i in range(0, len(text), max_chars)]

# ---------- UI ----------
uploaded_file = st.file_uploader(
    "Upload a file",
    type=["pdf", "docx", "txt", "csv", "xls", "xlsx", "pptx", "html", "htm", "png", "jpg", "jpeg"]
)

summary_type = st.selectbox(
    "Select summary type",
    [
        "Short Summary (Recommended)",
        "Topic-wise Summary",
        "Bullet-point Summary"
    ]
)

if uploaded_file and st.button("Generate Summary"):
    with st.spinner("Extracting text..."):
        text = extract_text(uploaded_file)

    if not text or not text.strip():
        st.error("Unsupported or empty file.")
    else:
        chunks = chunk_text(text)

        # ---------- MAP STEP ----------
        partial_summaries = []
        with st.spinner("Reading document..."):
            for chunk in chunks[:6]:
                res = client.chat.completions.create(
                    model=MODEL,
                    messages=[
                        {
                            "role": "system",
                            "content": "Extract only key ideas. No explanations."
                        },
                        {
                            "role": "user",
                            "content": chunk
                        }
                    ],
                    max_tokens=120
                )
                partial_summaries.append(res.choices[0].message.content.strip())

        combined_text = "\n".join(partial_summaries)

        # ---------- SUMMARY MODES ----------
        if summary_type == "Short Summary (Recommended)":
            final_prompt = """
Write a concise executive-style summary.

RULES:
- One short paragraph
- Simple language
- No headings
- No definitions
- Focus on overall meaning
- 120–150 tokens
"""

            max_tokens = 170

        elif summary_type == "Topic-wise Summary":
            final_prompt = """
Create a TOPIC-WISE summary.

RULES:
- Use clear topic headings
- 2–3 short sentences per topic
- Simple student-friendly language
- No formulas or definitions
- Do not repeat topics

FORMAT:
## Topic Name
Explanation
"""

            max_tokens = 600

        else:  # Bullet-point Summary
            final_prompt = """
Create a BULLET-POINT summary.

RULES:
- 8–12 bullet points
- Each bullet = one clear idea
- No explanations
- No definitions
- Easy to scan and act on
"""

            max_tokens = 300

        with st.spinner("Generating summary..."):
            final_response = client.chat.completions.create(
                model=MODEL,
                messages=[
                    {"role": "system", "content": final_prompt},
                    {"role": "user", "content": combined_text}
                ],
                max_tokens=max_tokens
            )

        st.success("Summary Generated")
        st.subheader("Summary")
        st.markdown(final_response.choices[0].message.content.strip())